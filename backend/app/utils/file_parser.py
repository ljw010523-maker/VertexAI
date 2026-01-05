"""
파일 파싱 유틸리티
다양한 파일 형식에서 텍스트를 추출합니다.
"""

import io
from typing import Optional
from pathlib import Path

# PDF
from pypdf import PdfReader

# Word
from docx import Document

# Excel
from openpyxl import load_workbook

# PowerPoint
from pptx import Presentation

# HTML
from bs4 import BeautifulSoup


class FileParser:
    """파일에서 텍스트를 추출하는 파서 클래스"""

    @staticmethod
    def parse_pdf(file_content: bytes) -> str:
        """PDF 파일에서 텍스트 추출"""
        try:
            pdf_file = io.BytesIO(file_content)
            reader = PdfReader(pdf_file)

            text_parts = []
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"[페이지 {page_num}]\n{text}")

            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"PDF 파싱 실패: {str(e)}")

    @staticmethod
    def parse_docx(file_content: bytes) -> str:
        """Word 문서에서 텍스트 추출"""
        try:
            docx_file = io.BytesIO(file_content)
            doc = Document(docx_file)

            text_parts = []

            # 문단 텍스트
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # 표 텍스트
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                if table_text:
                    text_parts.append("\n[표]\n" + "\n".join(table_text))

            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Word 문서 파싱 실패: {str(e)}")

    @staticmethod
    def parse_xlsx(file_content: bytes) -> str:
        """Excel 파일에서 텍스트 추출"""
        try:
            excel_file = io.BytesIO(file_content)
            workbook = load_workbook(excel_file, read_only=True, data_only=True)

            text_parts = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = [f"[시트: {sheet_name}]"]

                for row in sheet.iter_rows(values_only=True):
                    # None이 아닌 값들만 문자열로 변환
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        sheet_text.append(" | ".join(row_values))

                if len(sheet_text) > 1:  # 헤더만 있는게 아니라면
                    text_parts.append("\n".join(sheet_text))

            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Excel 파일 파싱 실패: {str(e)}")

    @staticmethod
    def parse_pptx(file_content: bytes) -> str:
        """PowerPoint 파일에서 텍스트 추출"""
        try:
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)

            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"[슬라이드 {slide_num}]"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if len(slide_text) > 1:
                    text_parts.append("\n".join(slide_text))

            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"PowerPoint 파일 파싱 실패: {str(e)}")

    @staticmethod
    def parse_html(file_content: bytes) -> str:
        """HTML 파일에서 텍스트 추출"""
        try:
            html_text = file_content.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html_text, 'html.parser')

            # 스크립트와 스타일 제거
            for script in soup(["script", "style"]):
                script.decompose()

            # 텍스트 추출
            text = soup.get_text()

            # 여러 줄바꿈을 하나로
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)

            return text
        except Exception as e:
            raise ValueError(f"HTML 파일 파싱 실패: {str(e)}")

    @staticmethod
    def parse_txt(file_content: bytes) -> str:
        """텍스트 파일에서 텍스트 추출"""
        try:
            # UTF-8로 시도
            try:
                return file_content.decode('utf-8')
            except UnicodeDecodeError:
                # 실패하면 다른 인코딩 시도
                for encoding in ['cp949', 'euc-kr', 'latin-1']:
                    try:
                        return file_content.decode(encoding)
                    except UnicodeDecodeError:
                        continue
                raise ValueError("지원하지 않는 텍스트 인코딩입니다")
        except Exception as e:
            raise ValueError(f"텍스트 파일 파싱 실패: {str(e)}")

    @classmethod
    def parse_file(cls, filename: str, file_content: bytes) -> str:
        """
        파일 확장자에 따라 적절한 파서를 선택하여 텍스트 추출

        Args:
            filename: 파일 이름
            file_content: 파일 바이트 내용

        Returns:
            추출된 텍스트

        Raises:
            ValueError: 지원하지 않는 파일 형식이거나 파싱 실패 시
        """
        # 파일 확장자 추출
        suffix = Path(filename).suffix.lower()

        # 확장자별 파서 매핑
        parsers = {
            '.pdf': cls.parse_pdf,
            '.docx': cls.parse_docx,
            '.doc': cls.parse_docx,  # .doc도 docx 파서로 시도
            '.xlsx': cls.parse_xlsx,
            '.xls': cls.parse_xlsx,
            '.pptx': cls.parse_pptx,
            '.ppt': cls.parse_pptx,
            '.html': cls.parse_html,
            '.htm': cls.parse_html,
            '.txt': cls.parse_txt,
            '.md': cls.parse_txt,
            '.csv': cls.parse_txt,
        }

        parser = parsers.get(suffix)
        if not parser:
            raise ValueError(
                f"지원하지 않는 파일 형식입니다: {suffix}\n"
                f"지원 형식: {', '.join(parsers.keys())}"
            )

        return parser(file_content)
